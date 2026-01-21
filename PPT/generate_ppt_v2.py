#!/usr/bin/env python3
"""
组会PPT生成脚本 v2
精简版：含图片支持
"""

import json
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# 配色方案
COLORS = {
    'primary_blue': RGBColor(0x14, 0x65, 0xC0),
    'dark_blue': RGBColor(0x0D, 0x47, 0xA1),
    'green': RGBColor(0x4C, 0xAF, 0x50),
    'dark_gray': RGBColor(0x33, 0x33, 0x33),
    'light_gray': RGBColor(0x75, 0x75, 0x75),
    'white': RGBColor(0xFF, 0xFF, 0xFF),
    'header_bg': RGBColor(0x14, 0x65, 0xC0),
}

# 幻灯片尺寸 (16:9)
SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)

def add_page_number(slide, page_num):
    """添加页码到右下角"""
    page_box = slide.shapes.add_textbox(
        SLIDE_WIDTH - Inches(0.8),
        SLIDE_HEIGHT - Inches(0.5),
        Inches(0.6),
        Inches(0.3)
    )
    tf = page_box.text_frame
    p = tf.paragraphs[0]
    p.text = str(page_num)
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = COLORS['primary_blue']
    p.alignment = PP_ALIGN.RIGHT

def add_header_bar(slide, title):
    """添加顶部标题栏"""
    header = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        SLIDE_WIDTH, Inches(0.9)
    )
    header.fill.solid()
    header.fill.fore_color.rgb = COLORS['header_bg']
    header.line.fill.background()

    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.2),
        SLIDE_WIDTH - Inches(1), Inches(0.6)
    )
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = COLORS['white']

def create_title_slide(prs, slide_data):
    """创建封面幻灯片"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(2.5),
        SLIDE_WIDTH - Inches(1), Inches(1.2)
    )
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = slide_data['title']
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = COLORS['dark_blue']
    p.alignment = PP_ALIGN.CENTER

    subtitle_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(3.8),
        SLIDE_WIDTH - Inches(1), Inches(0.6)
    )
    tf = subtitle_box.text_frame
    p = tf.paragraphs[0]
    p.text = slide_data.get('subtitle', '')
    p.font.size = Pt(24)
    p.font.color.rgb = COLORS['light_gray']
    p.alignment = PP_ALIGN.CENTER

    info_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(5.5),
        SLIDE_WIDTH - Inches(1), Inches(1)
    )
    tf = info_box.text_frame
    p = tf.paragraphs[0]
    p.text = f"{slide_data.get('presenter', '')}  |  {slide_data.get('date', '')}"
    p.font.size = Pt(18)
    p.font.color.rgb = COLORS['dark_gray']
    p.alignment = PP_ALIGN.CENTER

    if 'notes' in slide_data:
        slide.notes_slide.notes_text_frame.text = slide_data['notes']

    return slide

def create_image_grid_slide(prs, slide_data, page_num, base_path):
    """创建图片网格幻灯片（5张斑块照片）"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    add_header_bar(slide, slide_data['title'])

    images = slide_data.get('images', [])
    labels = slide_data.get('labels', [])

    # 5张图片布局: 上排3张，下排2张居中
    positions = [
        # 上排3张
        (Inches(0.5), Inches(1.2), Inches(4), Inches(2.8)),
        (Inches(4.67), Inches(1.2), Inches(4), Inches(2.8)),
        (Inches(8.83), Inches(1.2), Inches(4), Inches(2.8)),
        # 下排2张居中
        (Inches(2.58), Inches(4.2), Inches(4), Inches(2.8)),
        (Inches(6.75), Inches(4.2), Inches(4), Inches(2.8)),
    ]

    for i, (img_rel_path, label) in enumerate(zip(images, labels)):
        if i >= len(positions):
            break

        img_path = base_path / img_rel_path
        if not img_path.exists():
            continue

        left, top, width, height = positions[i]

        # 添加图片
        try:
            pic = slide.shapes.add_picture(str(img_path), left, top, width=width)
            # 调整高度保持比例
            if pic.height > height:
                ratio = height / pic.height
                pic.width = int(pic.width * ratio)
                pic.height = height
        except Exception as e:
            print(f"Warning: Could not add image {img_path}: {e}")
            continue

        # 添加标签
        label_box = slide.shapes.add_textbox(
            left, top + height + Inches(0.05),
            width, Inches(0.3)
        )
        tf = label_box.text_frame
        p = tf.paragraphs[0]
        p.text = label
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = COLORS['dark_gray']
        p.alignment = PP_ALIGN.CENTER

    add_page_number(slide, page_num)

    if 'notes' in slide_data:
        slide.notes_slide.notes_text_frame.text = slide_data['notes']

    return slide

def create_two_images_slide(prs, slide_data, page_num, base_path):
    """创建双图幻灯片（杀菌曲线对比）"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    add_header_bar(slide, slide_data['title'])

    images = slide_data.get('images', [])
    labels = slide_data.get('labels', [])

    # 左右两张图
    positions = [
        (Inches(0.3), Inches(1.1), Inches(6.3), Inches(5.8)),   # 左图
        (Inches(6.7), Inches(1.1), Inches(6.3), Inches(5.8)),   # 右图
    ]

    for i, (img_rel_path, label) in enumerate(zip(images, labels)):
        if i >= len(positions):
            break

        img_path = base_path / img_rel_path
        if not img_path.exists():
            print(f"Warning: Image not found: {img_path}")
            continue

        left, top, max_width, max_height = positions[i]

        try:
            pic = slide.shapes.add_picture(str(img_path), left, top, width=max_width)
            # 调整保持比例
            if pic.height > max_height:
                ratio = max_height / pic.height
                pic.width = int(pic.width * ratio)
                pic.height = max_height
        except Exception as e:
            print(f"Warning: Could not add image {img_path}: {e}")
            continue

        # 标签放在图片下方
        label_box = slide.shapes.add_textbox(
            left, top + Inches(5.0),
            max_width, Inches(0.4)
        )
        tf = label_box.text_frame
        p = tf.paragraphs[0]
        p.text = label
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = COLORS['dark_blue']
        p.alignment = PP_ALIGN.CENTER

    add_page_number(slide, page_num)

    if 'notes' in slide_data:
        slide.notes_slide.notes_text_frame.text = slide_data['notes']

    return slide

def create_summary_table_slide(prs, slide_data, page_num):
    """创建精简表格幻灯片"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    add_header_bar(slide, slide_data['title'])

    headers = slide_data.get('headers', [])
    rows = slide_data.get('rows', [])

    if not headers or not rows:
        add_page_number(slide, page_num)
        return slide

    num_cols = len(headers)
    num_rows = len(rows) + 1

    table_width = SLIDE_WIDTH - Inches(1.0)
    table_height = min(Inches(0.6) * num_rows, Inches(5.5))

    table = slide.shapes.add_table(
        num_rows, num_cols,
        Inches(0.5), Inches(1.2),
        table_width, table_height
    ).table

    # 设置列宽
    col_width = table_width / num_cols
    for i in range(num_cols):
        table.columns[i].width = int(col_width)

    # 填充表头
    for i, header in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = header
        cell.fill.solid()
        cell.fill.fore_color.rgb = COLORS['primary_blue']

        p = cell.text_frame.paragraphs[0]
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = COLORS['white']
        p.alignment = PP_ALIGN.CENTER
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE

    # 填充数据行
    for row_idx, row_data in enumerate(rows):
        for col_idx, cell_text in enumerate(row_data):
            cell = table.cell(row_idx + 1, col_idx)

            # 处理加粗标记
            text = str(cell_text)
            if '**' in text:
                text = text.replace('**', '')
                cell.text = text
                p = cell.text_frame.paragraphs[0]
                p.font.bold = True
                p.font.color.rgb = COLORS['green']
            else:
                cell.text = text
                p = cell.text_frame.paragraphs[0]
                p.font.color.rgb = COLORS['dark_gray']

            p.font.size = Pt(14)
            p.alignment = PP_ALIGN.CENTER
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE

            if row_idx % 2 == 1:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(0xF5, 0xF5, 0xF5)

    # 添加脚注
    if 'footnote' in slide_data:
        footnote_box = slide.shapes.add_textbox(
            Inches(0.5), SLIDE_HEIGHT - Inches(0.8),
            table_width, Inches(0.4)
        )
        tf = footnote_box.text_frame
        p = tf.paragraphs[0]
        p.text = slide_data['footnote']
        p.font.size = Pt(12)
        p.font.italic = True
        p.font.color.rgb = COLORS['light_gray']

    add_page_number(slide, page_num)

    if 'notes' in slide_data:
        slide.notes_slide.notes_text_frame.text = slide_data['notes']

    return slide

def create_conclusion_slide(prs, slide_data, page_num):
    """创建结论幻灯片"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    add_header_bar(slide, slide_data['title'])

    # 左栏：主要发现
    left_title = slide.shapes.add_textbox(
        Inches(0.5), Inches(1.2),
        Inches(5.5), Inches(0.5)
    )
    tf = left_title.text_frame
    p = tf.paragraphs[0]
    p.text = "Main Findings" if 'findings' in slide_data else "主要发现"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = COLORS['primary_blue']

    left_content = slide.shapes.add_textbox(
        Inches(0.5), Inches(1.8),
        Inches(5.8), Inches(4.5)
    )
    tf = left_content.text_frame
    tf.word_wrap = True

    for i, item in enumerate(slide_data.get('findings', [])):
        if i > 0:
            p = tf.add_paragraph()
        else:
            p = tf.paragraphs[0]
        p.space_before = Pt(8)

        # 处理加粗
        text = f"• {item}"
        if '**' in text:
            parts = text.split('**')
            for j, part in enumerate(parts):
                run = p.add_run()
                run.text = part
                if j % 2 == 1:
                    run.font.bold = True
                    run.font.color.rgb = COLORS['green']
                else:
                    run.font.color.rgb = COLORS['dark_gray']
                run.font.size = Pt(16)
        else:
            p.text = text
            p.font.size = Pt(16)
            p.font.color.rgb = COLORS['dark_gray']

    # 右栏：下一步计划
    right_title = slide.shapes.add_textbox(
        Inches(6.8), Inches(1.2),
        Inches(5.5), Inches(0.5)
    )
    tf = right_title.text_frame
    p = tf.paragraphs[0]
    p.text = "Next Steps" if 'next_steps' in slide_data else "下一步计划"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = COLORS['green']

    right_content = slide.shapes.add_textbox(
        Inches(6.8), Inches(1.8),
        Inches(5.8), Inches(4.5)
    )
    tf = right_content.text_frame
    tf.word_wrap = True

    for i, item in enumerate(slide_data.get('next_steps', [])):
        if i > 0:
            p = tf.add_paragraph()
        else:
            p = tf.paragraphs[0]
        p.space_before = Pt(8)
        p.text = f"→ {item}"
        p.font.size = Pt(16)
        p.font.color.rgb = COLORS['dark_gray']

    add_page_number(slide, page_num)

    if 'notes' in slide_data:
        slide.notes_slide.notes_text_frame.text = slide_data['notes']

    return slide

def create_thank_you_slide(prs, slide_data, page_num):
    """创建致谢幻灯片"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(2.8),
        SLIDE_WIDTH - Inches(1), Inches(1.2)
    )
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Thank You!"
    p.font.size = Pt(54)
    p.font.bold = True
    p.font.color.rgb = COLORS['primary_blue']
    p.alignment = PP_ALIGN.CENTER

    qa_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(4.2),
        SLIDE_WIDTH - Inches(1), Inches(0.6)
    )
    tf = qa_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Questions & Discussion"
    p.font.size = Pt(24)
    p.font.color.rgb = COLORS['light_gray']
    p.alignment = PP_ALIGN.CENTER

    add_page_number(slide, page_num)

    if 'notes' in slide_data:
        slide.notes_slide.notes_text_frame.text = slide_data['notes']

    return slide

def generate_ppt(outline_path, output_path, base_path):
    """根据JSON框架生成PPT"""
    with open(outline_path, 'r', encoding='utf-8') as f:
        outline = json.load(f)

    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    page_num = 0

    for slide_data in outline.get('slides', []):
        page_num += 1
        slide_type = slide_data.get('type', 'content')

        if slide_type == 'title':
            create_title_slide(prs, slide_data)
        elif slide_type == 'image_grid':
            create_image_grid_slide(prs, slide_data, page_num, base_path)
        elif slide_type == 'two_images':
            create_two_images_slide(prs, slide_data, page_num, base_path)
        elif slide_type == 'table':
            create_summary_table_slide(prs, slide_data, page_num)
        elif slide_type == 'conclusion':
            create_conclusion_slide(prs, slide_data, page_num)
        elif slide_type == 'thank_you':
            create_thank_you_slide(prs, slide_data, page_num)
        else:
            create_summary_table_slide(prs, slide_data, page_num)

    prs.save(output_path)
    print(f"已生成: {output_path}")
    return page_num

def main():
    base_dir = Path(r"C:\Users\36094\Desktop\EcAZPhageDocumentation\PPT\20260126")
    project_root = Path(r"C:\Users\36094\Desktop\EcAZPhageDocumentation")

    # 生成中文版
    cn_outline = base_dir / "ppt_outline_cn_v2.json"
    cn_output = base_dir / "组会汇报_v2.pptx"

    if cn_outline.exists():
        pages = generate_ppt(cn_outline, cn_output, project_root)
        print(f"  共 {pages} 页")

    # 生成英文版
    en_outline = base_dir / "ppt_outline_en_v2.json"
    en_output = base_dir / "GroupMeeting_v2.pptx"

    if en_outline.exists():
        pages = generate_ppt(en_outline, en_output, project_root)
        print(f"  共 {pages} 页")

    print("\n生成完成！")

if __name__ == "__main__":
    main()
