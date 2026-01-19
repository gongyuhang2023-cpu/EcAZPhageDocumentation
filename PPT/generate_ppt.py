#!/usr/bin/env python3
"""
组会PPT生成脚本
根据JSON框架生成专业的PPT文件
"""

import json
import re
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# 配色方案
COLORS = {
    'primary_blue': RGBColor(0x14, 0x65, 0xC0),
    'dark_blue': RGBColor(0x0D, 0x47, 0xA1),
    'green': RGBColor(0x4C, 0xAF, 0x50),
    'dark_gray': RGBColor(0x33, 0x33, 0x33),
    'light_gray': RGBColor(0x75, 0x75, 0x75),
    'white': RGBColor(0xFF, 0xFF, 0xFF),
    'header_bg': RGBColor(0x14, 0x65, 0xC0),
}

# 幻灯片尺寸 (16:9)
SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)

def add_page_number(slide, page_num):
    """添加页码到右下角"""
    page_box = slide.shapes.add_textbox(
        SLIDE_WIDTH - Inches(0.8),
        SLIDE_HEIGHT - Inches(0.5),
        Inches(0.6),
        Inches(0.3)
    )
    tf = page_box.text_frame
    p = tf.paragraphs[0]
    p.text = str(page_num)
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = COLORS['primary_blue']
    p.alignment = PP_ALIGN.RIGHT

def add_header_bar(slide, title):
    """添加顶部标题栏"""
    # 蓝色背景条
    header = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        SLIDE_WIDTH, Inches(0.9)
    )
    header.fill.solid()
    header.fill.fore_color.rgb = COLORS['header_bg']
    header.line.fill.background()

    # 标题文字
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.2),
        SLIDE_WIDTH - Inches(1), Inches(0.6)
    )
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = COLORS['white']

def format_text_with_bold(paragraph, text):
    """处理**text**格式的加粗文本"""
    # 清空现有内容
    paragraph.clear()

    # 使用正则分割文本
    parts = re.split(r'(\*\*.*?\*\*)', text)

    for part in parts:
        if part.startswith('**') and part.endswith('**'):
            # 加粗文本
            run = paragraph.add_run()
            run.text = part[2:-2]
            run.font.bold = True
            run.font.color.rgb = COLORS['green']
        else:
            # 普通文本
            if part:
                run = paragraph.add_run()
                run.text = part

def create_title_slide(prs, slide_data):
    """创建封面幻灯片"""
    slide_layout = prs.slide_layouts[6]  # 空白布局
    slide = prs.slides.add_slide(slide_layout)

    # 标题
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(2.5),
        SLIDE_WIDTH - Inches(1), Inches(1.2)
    )
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = slide_data['title']
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = COLORS['dark_blue']
    p.alignment = PP_ALIGN.CENTER

    # 副标题
    subtitle_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(3.8),
        SLIDE_WIDTH - Inches(1), Inches(0.6)
    )
    tf = subtitle_box.text_frame
    p = tf.paragraphs[0]
    p.text = slide_data.get('subtitle', '')
    p.font.size = Pt(24)
    p.font.color.rgb = COLORS['light_gray']
    p.alignment = PP_ALIGN.CENTER

    # 演讲者和日期
    info_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(5.5),
        SLIDE_WIDTH - Inches(1), Inches(1)
    )
    tf = info_box.text_frame
    p = tf.paragraphs[0]
    p.text = f"{slide_data.get('presenter', '')}  |  {slide_data.get('date', '')}"
    p.font.size = Pt(18)
    p.font.color.rgb = COLORS['dark_gray']
    p.alignment = PP_ALIGN.CENTER

    # 添加备注
    if 'notes' in slide_data:
        slide.notes_slide.notes_text_frame.text = slide_data['notes']

    return slide

def create_content_slide(prs, slide_data, page_num):
    """创建内容幻灯片（要点列表）"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    add_header_bar(slide, slide_data['title'])

    # 内容区域
    content_top = Inches(1.3)
    content_box = slide.shapes.add_textbox(
        Inches(0.8), content_top,
        SLIDE_WIDTH - Inches(1.6), SLIDE_HEIGHT - content_top - Inches(0.8)
    )
    tf = content_box.text_frame
    tf.word_wrap = True

    for i, item in enumerate(slide_data.get('content', [])):
        if i > 0:
            p = tf.add_paragraph()
        else:
            p = tf.paragraphs[0]

        p.space_before = Pt(12)
        p.space_after = Pt(8)

        # 添加bullet
        format_text_with_bold(p, f"• {item}")
        p.font.size = Pt(22)
        p.font.color.rgb = COLORS['dark_gray']

    add_page_number(slide, page_num)

    if 'notes' in slide_data:
        slide.notes_slide.notes_text_frame.text = slide_data['notes']

    return slide

def create_table_slide(prs, slide_data, page_num):
    """创建表格幻灯片"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    add_header_bar(slide, slide_data['title'])

    headers = slide_data.get('headers', [])
    rows = slide_data.get('rows', [])

    if not headers or not rows:
        add_page_number(slide, page_num)
        return slide

    # 计算表格尺寸
    num_cols = len(headers)
    num_rows = len(rows) + 1  # +1 for header

    table_width = SLIDE_WIDTH - Inches(1.6)
    table_height = min(Inches(0.5) * num_rows, Inches(4.5))

    # 创建表格
    table = slide.shapes.add_table(
        num_rows, num_cols,
        Inches(0.8), Inches(1.3),
        table_width, table_height
    ).table

    # 设置列宽
    col_width = table_width / num_cols
    for i in range(num_cols):
        table.columns[i].width = int(col_width)

    # 填充表头
    for i, header in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = header
        cell.fill.solid()
        cell.fill.fore_color.rgb = COLORS['primary_blue']

        p = cell.text_frame.paragraphs[0]
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = COLORS['white']
        p.alignment = PP_ALIGN.CENTER
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE

    # 填充数据行
    for row_idx, row_data in enumerate(rows):
        for col_idx, cell_text in enumerate(row_data):
            cell = table.cell(row_idx + 1, col_idx)

            p = cell.text_frame.paragraphs[0]
            format_text_with_bold(p, str(cell_text))
            p.font.size = Pt(13)
            p.alignment = PP_ALIGN.CENTER
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE

            # 交替行背景色
            if row_idx % 2 == 1:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(0xF5, 0xF5, 0xF5)

    # 添加脚注
    if 'footnote' in slide_data:
        footnote_box = slide.shapes.add_textbox(
            Inches(0.8), SLIDE_HEIGHT - Inches(1),
            table_width, Inches(0.4)
        )
        tf = footnote_box.text_frame
        p = tf.paragraphs[0]
        p.text = slide_data['footnote']
        p.font.size = Pt(12)
        p.font.italic = True
        p.font.color.rgb = COLORS['light_gray']

    add_page_number(slide, page_num)

    if 'notes' in slide_data:
        slide.notes_slide.notes_text_frame.text = slide_data['notes']

    return slide

def create_conclusion_slide(prs, slide_data, page_num):
    """创建结论幻灯片（双栏布局）"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    add_header_bar(slide, slide_data['title'])

    # 左栏：主要发现
    left_title = slide.shapes.add_textbox(
        Inches(0.5), Inches(1.2),
        Inches(5.5), Inches(0.5)
    )
    tf = left_title.text_frame
    p = tf.paragraphs[0]
    p.text = "Main Findings" if 'findings' in slide_data else "主要发现"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = COLORS['primary_blue']

    left_content = slide.shapes.add_textbox(
        Inches(0.5), Inches(1.8),
        Inches(5.8), Inches(4.5)
    )
    tf = left_content.text_frame
    tf.word_wrap = True

    for i, item in enumerate(slide_data.get('findings', [])):
        if i > 0:
            p = tf.add_paragraph()
        else:
            p = tf.paragraphs[0]
        p.space_before = Pt(8)
        format_text_with_bold(p, f"• {item}")
        p.font.size = Pt(18)
        p.font.color.rgb = COLORS['dark_gray']

    # 右栏：下一步计划
    right_title = slide.shapes.add_textbox(
        Inches(6.8), Inches(1.2),
        Inches(5.5), Inches(0.5)
    )
    tf = right_title.text_frame
    p = tf.paragraphs[0]
    p.text = "Next Steps" if 'next_steps' in slide_data else "下一步计划"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = COLORS['green']

    right_content = slide.shapes.add_textbox(
        Inches(6.8), Inches(1.8),
        Inches(5.8), Inches(4.5)
    )
    tf = right_content.text_frame
    tf.word_wrap = True

    for i, item in enumerate(slide_data.get('next_steps', [])):
        if i > 0:
            p = tf.add_paragraph()
        else:
            p = tf.paragraphs[0]
        p.space_before = Pt(8)
        p.text = f"→ {item}"
        p.font.size = Pt(18)
        p.font.color.rgb = COLORS['dark_gray']

    add_page_number(slide, page_num)

    if 'notes' in slide_data:
        slide.notes_slide.notes_text_frame.text = slide_data['notes']

    return slide

def create_thank_you_slide(prs, slide_data, page_num):
    """创建致谢幻灯片"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # 大标题
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(2.8),
        SLIDE_WIDTH - Inches(1), Inches(1.2)
    )
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Thank You!"
    p.font.size = Pt(54)
    p.font.bold = True
    p.font.color.rgb = COLORS['primary_blue']
    p.alignment = PP_ALIGN.CENTER

    # Q&A
    qa_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(4.2),
        SLIDE_WIDTH - Inches(1), Inches(0.6)
    )
    tf = qa_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Questions & Discussion"
    p.font.size = Pt(24)
    p.font.color.rgb = COLORS['light_gray']
    p.alignment = PP_ALIGN.CENTER

    add_page_number(slide, page_num)

    if 'notes' in slide_data:
        slide.notes_slide.notes_text_frame.text = slide_data['notes']

    return slide

def generate_ppt(outline_path, output_path):
    """根据JSON框架生成PPT"""
    # 读取框架
    with open(outline_path, 'r', encoding='utf-8') as f:
        outline = json.load(f)

    # 创建演示文稿
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    page_num = 0

    for slide_data in outline.get('slides', []):
        page_num += 1
        slide_type = slide_data.get('type', 'content')

        if slide_type == 'title':
            create_title_slide(prs, slide_data)
        elif slide_type == 'content':
            create_content_slide(prs, slide_data, page_num)
        elif slide_type == 'table':
            create_table_slide(prs, slide_data, page_num)
        elif slide_type == 'two_column':
            create_content_slide(prs, slide_data, page_num)  # 简化处理
        elif slide_type == 'conclusion':
            create_conclusion_slide(prs, slide_data, page_num)
        elif slide_type == 'thank_you':
            create_thank_you_slide(prs, slide_data, page_num)
        else:
            create_content_slide(prs, slide_data, page_num)

    # 保存
    prs.save(output_path)
    print(f"已生成: {output_path}")
    return page_num

def main():
    import sys

    # 默认路径
    base_dir = Path(r"C:\Users\36094\Desktop\EcAZPhageDocumentation\PPT\20260126")

    # 生成中文版
    cn_outline = base_dir / "ppt_outline_cn.json"
    cn_output = base_dir / "组会汇报_v1.pptx"

    if cn_outline.exists():
        pages = generate_ppt(cn_outline, cn_output)
        print(f"  共 {pages} 页")

    # 生成英文版
    en_outline = base_dir / "ppt_outline_en.json"
    en_output = base_dir / "GroupMeeting_v1.pptx"

    if en_outline.exists():
        pages = generate_ppt(en_outline, en_output)
        print(f"  共 {pages} 页")

    print("\n生成完成！")

if __name__ == "__main__":
    main()
